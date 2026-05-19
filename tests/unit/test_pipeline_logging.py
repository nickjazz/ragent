"""T2v.42 — Per-step structured logging for ingest pipeline components.

Each pipeline component emits ``ingest.step.{started,ok,failed}`` via
``structlog.get_logger("ragent.ingest")``. Failures map to the existing
error-code taxonomy (``PIPELINE_UNROUTABLE`` / ``EMBEDDER_ERROR`` /
``ES_WRITE_ERROR`` / ``PIPELINE_TIMEOUT``). The worker emits terminal
``ingest.ready`` / ``ingest.failed`` events with totals.
"""

from __future__ import annotations

import contextlib

import structlog

from ragent.pipelines.observability import (
    IngestStepError,
    bind_ingest_context,
    log_ingest_step,
    wrap_pipeline_component,
)


def _wrap(component, *, step, error_code="PIPELINE_UNEXPECTED_ERROR"):
    """Test-local ingest-namespaced shim; mirrors the dropped wrap_component_run alias."""
    return wrap_pipeline_component(component, namespace="ingest", step=step, error_code=error_code)


# ---------------------------------------------------------------------------
# wrap_pipeline_component (ingest namespace) — happy path
# ---------------------------------------------------------------------------


class _FakeComponent:
    def __init__(self) -> None:
        self.called_with: dict | None = None

    def run(self, documents: list, **kwargs) -> dict:
        self.called_with = {"documents": documents, **kwargs}
        return {"documents": [{"out": True} for _ in range(2)]}


def test_wrap_emits_started_and_ok_with_expected_fields() -> None:
    comp = _FakeComponent()
    _wrap(comp, step="embedder")
    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-1", mime_type="text/markdown"),
    ):
        out = comp.run(documents=[{"a": 1}, {"a": 2}, {"a": 3}])
    assert out == {"documents": [{"out": True}, {"out": True}]}

    events = [e for e in logs if e.get("event", "").startswith("ingest.step.")]
    assert [e["event"] for e in events] == ["ingest.step.started", "ingest.step.ok"]
    started, ok = events
    assert started["step"] == "embedder"
    assert started["document_id"] == "DOC-1"
    assert started["mime_type"] == "text/markdown"
    assert ok["step"] == "embedder"
    assert ok["document_id"] == "DOC-1"
    assert ok["mime_type"] == "text/markdown"
    assert isinstance(ok["duration_ms"], int)
    assert ok["duration_ms"] >= 0
    assert ok["atoms_in"] == 3
    assert ok["chunks_out"] == 2


# ---------------------------------------------------------------------------
# wrap_component_run — failure path
# ---------------------------------------------------------------------------


class _BoomComponent:
    def run(self, documents: list) -> dict:
        raise RuntimeError("boom")


def test_wrap_emits_failed_with_error_code_and_reraises() -> None:
    comp = _BoomComponent()
    _wrap(comp, step="embedder", error_code="EMBEDDER_ERROR")
    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-2", mime_type="text/plain"),
    ):
        try:
            comp.run(documents=[1, 2])
        except RuntimeError as exc:
            assert str(exc) == "boom"
        else:  # pragma: no cover
            raise AssertionError("expected RuntimeError")

    failed = [e for e in logs if e.get("event") == "ingest.step.failed"]
    assert len(failed) == 1
    e = failed[0]
    assert e["step"] == "embedder"
    assert e["document_id"] == "DOC-2"
    assert e["mime_type"] == "text/plain"
    assert e["error_code"] == "EMBEDDER_ERROR"
    assert "boom" in e["error"]
    assert isinstance(e["duration_ms"], int)


def test_wrap_failure_with_explicit_error_code_via_exception() -> None:
    """Components that raise ``IngestStepError`` carry their own error_code."""

    class _RouterRaiser:
        def run(self, documents: list) -> dict:
            raise IngestStepError("unroutable mime", error_code="PIPELINE_UNROUTABLE")

    comp = _RouterRaiser()
    _wrap(comp, step="router", error_code="PIPELINE_UNROUTABLE")
    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-3", mime_type="text/x-bogus"),
        contextlib.suppress(IngestStepError),
    ):
        comp.run(documents=[])

    failed = [e for e in logs if e.get("event") == "ingest.step.failed"]
    assert failed[0]["error_code"] == "PIPELINE_UNROUTABLE"


# ---------------------------------------------------------------------------
# Pipeline-level ordering: build_ingest_pipeline wraps each step
# ---------------------------------------------------------------------------


def test_build_ingest_pipeline_wraps_steps_in_order(monkeypatch) -> None:
    """Run the v2 pipeline end-to-end with mocks; assert step events emitted in order."""
    from unittest.mock import MagicMock

    from haystack_integrations.document_stores.elasticsearch import ElasticsearchDocumentStore

    from ragent.pipelines.factory import DocumentEmbedder, build_ingest_pipeline

    class _StubEmbedder:
        def embed(self, texts):
            return [[0.1] * 4 for _ in texts]

    embedder = DocumentEmbedder(_StubEmbedder())
    document_store = MagicMock(spec=ElasticsearchDocumentStore)
    document_store.write_documents.return_value = 0

    pipe = build_ingest_pipeline(embedder=embedder, document_store=document_store)

    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-PIPE", mime_type="text/plain"),
    ):
        pipe.run(
            {
                "loader": {
                    "content": "hello world. this is a sentence. and another.",
                    "mime_type": "text/plain",
                    "document_id": "DOC-PIPE",
                }
            }
        )

    step_events = [e for e in logs if e.get("event", "").startswith("ingest.step.")]
    pairs: list[tuple[str, str]] = []
    for ev in step_events:
        pairs.append((ev["event"].split(".")[-1], ev["step"]))
    started_steps = [s for k, s in pairs if k == "started"]
    ok_steps = [s for k, s in pairs if k == "ok"]
    # v2 graph: load → split → chunker → embedder → writer.
    expected = ["load", "split", "chunker", "embedder", "writer"]
    assert started_steps == expected
    assert ok_steps == expected
    for ev in step_events:
        assert ev["document_id"] == "DOC-PIPE"
        assert ev["mime_type"] == "text/plain"


# ---------------------------------------------------------------------------
# log_ingest_step terminal helpers
# ---------------------------------------------------------------------------


def test_wrap_writer_chunks_out_uses_int_documents_written() -> None:
    """Haystack DocumentWriter.run returns {"documents_written": int}."""

    class _IntWriter:
        def run(self, documents: list) -> dict:
            return {"documents_written": len(documents)}

    comp = _IntWriter()
    _wrap(comp, step="writer", error_code="ES_WRITE_ERROR")
    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-W", mime_type="text/plain"),
    ):
        comp.run(documents=[1, 2, 3, 4])

    ok = [e for e in logs if e.get("event") == "ingest.step.ok"][0]
    assert ok["chunks_out"] == 4
    assert ok["atoms_in"] == 4


def test_log_ingest_step_ready_emits_terminal_event() -> None:
    with structlog.testing.capture_logs() as logs:
        log_ingest_step.ready(document_id="DOC-9", chunks_total=7, duration_ms_total=42)
    e = [x for x in logs if x.get("event") == "ingest.ready"][0]
    assert e["document_id"] == "DOC-9"
    assert e["chunks_total"] == 7
    assert e["duration_ms_total"] == 42


def test_log_ingest_step_failed_emits_terminal_event() -> None:
    with structlog.testing.capture_logs() as logs:
        log_ingest_step.failed(
            document_id="DOC-9", reason="pipeline_error", error_code="EMBEDDER_ERROR"
        )
    e = [x for x in logs if x.get("event") == "ingest.failed"][0]
    assert e["document_id"] == "DOC-9"
    assert e["reason"] == "pipeline_error"
    assert e["error_code"] == "EMBEDDER_ERROR"


def test_extra_context_vars_appear_in_step_logs() -> None:
    """Extra vars bound to structlog context (e.g. file_size_bytes) propagate to step logs.

    _ctx() must pass through all bound contextvars, not just document_id and mime_type.
    This is the mechanism that puts file_size_bytes in the load step and splitter in
    the split step.
    """
    comp = _FakeComponent()
    _wrap(comp, step="load")

    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-X", mime_type="text/plain"),
    ):
        structlog.contextvars.bind_contextvars(file_size_bytes=1234)
        comp.run(documents=[{"a": 1}])

    ok = [e for e in logs if e.get("event") == "ingest.step.ok"][0]
    assert ok["file_size_bytes"] == 1234


def test_splitter_context_var_appears_in_split_step_log() -> None:
    """splitter bound to context during _MimeAwareSplitter.run() appears in ingest.step.ok.

    _MimeAwareSplitter must bind splitter=<name> to structlog context so the
    split step's ingest.step.ok log carries the routing decision.
    """
    import io

    from pptx import Presentation
    from pptx.util import Inches

    from ragent.pipelines.factory import _MimeAwareSplitter
    from ragent.schemas.ingest import IngestMime

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = "hi"
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()

    from haystack.dataclasses import Document

    splitter = _MimeAwareSplitter()
    _wrap(splitter, step="split")
    with (
        structlog.testing.capture_logs() as logs,
        bind_ingest_context(document_id="DOC-Y", mime_type=IngestMime.PPTX),
    ):
        doc = Document(content=None, meta={"mime_type": IngestMime.PPTX, "raw_bytes": data})
        splitter.run(documents=[doc])

    ok = [e for e in logs if e.get("event") == "ingest.step.ok"][0]
    assert ok["splitter"] == "pptx"
