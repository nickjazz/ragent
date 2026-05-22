"""T2v.38 — Mime-aware splitter dispatches per meta['mime_type'].

Unknown MIME raises IngestStepError(error_code=PIPELINE_UNROUTABLE).
"""

from __future__ import annotations

import pytest
from haystack.dataclasses import Document

from ragent.pipelines.ingest import _MimeAwareSplitter
from ragent.pipelines.observability import IngestStepError


def test_plain_routes_to_document_splitter() -> None:
    out = _MimeAwareSplitter().run(
        [Document(content="hello world.", meta={"mime_type": "text/plain"})]
    )["documents"]
    assert len(out) >= 1
    for a in out:
        assert a.meta.get("raw_content")  # plain fallback uses content as raw


def test_markdown_routes_to_markdown_ast_splitter() -> None:
    out = _MimeAwareSplitter().run(
        [Document(content="# H\n\np\n\n```\nx=1\n```", meta={"mime_type": "text/markdown"})]
    )["documents"]
    raws = "".join(a.meta.get("raw_content", "") for a in out)
    assert "```" in raws


def test_html_routes_to_html_ast_splitter() -> None:
    out = _MimeAwareSplitter().run(
        [Document(content="<p>hello</p><script>x=1</script>", meta={"mime_type": "text/html"})]
    )["documents"]
    raws = " ".join(a.meta.get("raw_content", "") for a in out)
    assert "x=1" not in raws


def test_docx_routes_to_docx_ast_splitter(tmp_path) -> None:
    import io

    from docx import Document as DocxDocument

    doc = DocxDocument()
    for p in list(doc.paragraphs):
        p._element.getparent().remove(p._element)
    doc.add_paragraph("DOCX content here.")
    buf = io.BytesIO()
    doc.save(buf)
    data = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    out = _MimeAwareSplitter().run(
        [Document(content=None, meta={"mime_type": mime, "raw_bytes": data})]
    )["documents"]
    assert len(out) >= 1
    assert any("DOCX content here." in (a.content or "") for a in out)


def test_pptx_routes_to_pptx_ast_splitter() -> None:
    import io

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    txBox.text_frame.text = "PPTX slide text"
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    out = _MimeAwareSplitter().run(
        [Document(content=None, meta={"mime_type": mime, "raw_bytes": data})]
    )["documents"]
    assert len(out) == 1
    assert "PPTX slide text" in out[0].content


def test_unknown_mime_raises_pipeline_unroutable() -> None:
    with pytest.raises(IngestStepError) as exc:
        _MimeAwareSplitter().run([Document(content="x", meta={"mime_type": "image/png"})])
    assert exc.value.error_code == "PIPELINE_UNROUTABLE"


def test_pdf_routes_to_pdf_ast_splitter() -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(False)
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, text="PDF route test")
    data = bytes(pdf.output())

    out = _MimeAwareSplitter().run(
        [Document(content=None, meta={"mime_type": "application/pdf", "raw_bytes": data})]
    )["documents"]
    assert len(out) == 1
    assert "PDF route test" in out[0].content
