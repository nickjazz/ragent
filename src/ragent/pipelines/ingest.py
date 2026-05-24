"""C4 / T2v.30-T2v.39 — v2 ingest pipeline factory.

Graph: ``_TextLoader → _MimeAwareSplitter → _BudgetChunker →
DocumentEmbedder`` (ES only; embedder bulk-writes directly).

Splitter dispatches per ``meta["mime_type"]``:
- ``text/plain``      → Haystack ``DocumentSplitter`` (passage)
- ``text/markdown``   → ``_MarkdownASTSplitter`` (mistletoe)
- ``text/html``       → ``_HtmlASTSplitter`` (selectolax)
- ``application/pdf`` → ``_PdfASTSplitter`` (pymupdf4llm to_markdown + RapidOCR auto)

Each splitter emits atoms whose ``meta["raw_content"]`` is the original
byte slice (markdown fences / HTML fragments preserved). ``_BudgetChunker``
is mime-agnostic: greedy-pack to ``CHUNK_TARGET_CHARS``, hard-split atoms
exceeding ``CHUNK_MAX_CHARS`` with ``CHUNK_OVERLAP_CHARS`` overlap.
"""

from __future__ import annotations

import dataclasses
import io
import re
from concurrent.futures import ThreadPoolExecutor
from typing import Any

import pymupdf4llm
import structlog
from haystack.components.preprocessors import DocumentSplitter
from haystack.core.component import component
from haystack.core.pipeline import Pipeline
from haystack.dataclasses import Document

from ragent.errors.codes import TaskErrorCode
from ragent.pipelines.observability import IngestStepError, wrap_pipeline_component
from ragent.schemas.ingest import IngestMime
from ragent.security.archive_guard import INGEST_MAX_PDF_PAGES
from ragent.utility.env import float_env, int_env

_logger = structlog.get_logger(__name__)

# Single mime-agnostic budget profile (replaces v1 EN/CJK/CSV constants).
CHUNK_TARGET_CHARS = int_env("CHUNK_TARGET_CHARS", 1000)
CHUNK_MAX_CHARS = int_env("CHUNK_MAX_CHARS", 1500)
CHUNK_OVERLAP_CHARS = int_env("CHUNK_OVERLAP_CHARS", 100)
# Per-atom hard-split safety cap: an atom that exceeds this many pieces is
# treated as misconfiguration (overlap ≥ target produces tiny advance steps
# and can blow up to millions of chunks on a 1 MB atom).
CHUNK_MAX_PIECES_PER_ATOM = int_env("CHUNK_MAX_PIECES_PER_ATOM", 10_000)
# PDF page margin in points (1 pt ≈ 0.35 mm). Header/footer zones at the
# top and bottom of each page are excluded from extraction when > 0.
INGEST_PDF_MARGIN_PTS = float_env("INGEST_PDF_MARGIN_PTS", 0.0)
# PPTX placeholder types to exclude (header=14, footer=15, date=16, slide_number=13).
# Integer values used so python-pptx stays a lazy import inside run().
_PPTX_SKIP_PH: frozenset[int] = frozenset({13, 14, 15, 16})


def validate_chunk_config() -> None:
    """Bootstrap-time invariant check.

    Called from ``bootstrap.guard.enforce`` so a misconfigured env aborts
    process boot cleanly rather than crashing every importer (tests,
    linters) with a module-level RuntimeError.
    """
    if CHUNK_TARGET_CHARS <= CHUNK_OVERLAP_CHARS:
        raise RuntimeError(
            "CHUNK_TARGET_CHARS must be > CHUNK_OVERLAP_CHARS; "
            f"got target={CHUNK_TARGET_CHARS}, overlap={CHUNK_OVERLAP_CHARS}"
        )


ALLOWED_MIMES = (
    "text/plain",
    "text/markdown",
    "text/html",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    IngestMime.PDF,
)

# ---------------------------------------------------------------------------
# _TextLoader (T2v.30/31)
# ---------------------------------------------------------------------------


@component
class _TextLoader:
    """Build a single ``Document`` from inline content + per-document meta.

    The worker calls ``run(content=..., mime_type=..., document_id=...)`` so
    the loader replaces v1's ``TextFileToDocument`` + tempfile dance.
    """

    @component.output_types(documents=list[Document])
    def run(
        self,
        content: str,
        mime_type: str,
        document_id: str | None = None,
        source_url: str | None = None,
        source_title: str | None = None,
        source_app: str | None = None,
        source_meta: str | None = None,
        content_bytes: bytes | None = None,
    ) -> dict:
        meta: dict[str, Any] = {"mime_type": mime_type}
        for k, v in (
            ("document_id", document_id),
            ("source_url", source_url),
            ("source_title", source_title),
            ("source_app", source_app),
            ("source_meta", source_meta),
        ):
            if v is not None:
                meta[k] = v
        if content_bytes is not None:
            meta["raw_bytes"] = content_bytes
        return {"documents": [Document(content=content, meta=meta)]}


# ---------------------------------------------------------------------------
# _MarkdownASTSplitter (T2v.32/33)
# ---------------------------------------------------------------------------


_MD_BLOCK_TYPES = (
    "Heading",
    "Paragraph",
    "CodeFence",
    "List",
    "Table",
    "Quote",
    "ThematicBreak",
    "HtmlBlock",
)

# Strip common markdown markers from a rendered block so the `content`
# field embeds + BM25-indexes prose text rather than syntax noise. Fenced
# code (``` markers) and inline code spans (`x`) are unwrapped to their
# inner text. Heading hashes, bullet/quote prefixes, and emphasis markers
# are dropped. The original markup is always preserved in `raw_content`.
_MD_FENCE_RE = re.compile(r"^```[^\n]*\n(.*?)(?:\n```\s*)?$", re.DOTALL)
_MD_INLINE_CODE_RE = re.compile(r"`+([^`]+)`+")
_MD_HEADING_RE = re.compile(r"^#{1,6}\s+", re.MULTILINE)
_MD_LIST_PREFIX_RE = re.compile(r"^\s*(?:[-*+]|\d+\.)\s+", re.MULTILINE)
_MD_QUOTE_PREFIX_RE = re.compile(r"^>\s?", re.MULTILINE)
_MD_EMPHASIS_RE = re.compile(r"(\*{1,3}|_{1,3})(\S(?:.*?\S)?)\1")


def _md_plain(raw: str, type_name: str) -> str:
    """Reduce a rendered markdown block to its prose content."""
    if type_name == "CodeFence":
        m = _MD_FENCE_RE.match(raw.rstrip())
        return (m.group(1) if m else raw).rstrip()
    text = raw
    if type_name == "Heading":
        text = _MD_HEADING_RE.sub("", text)
    if type_name in ("List", "Quote"):
        text = _MD_LIST_PREFIX_RE.sub("", text)
        text = _MD_QUOTE_PREFIX_RE.sub("", text)
    text = _MD_INLINE_CODE_RE.sub(r"\1", text)
    text = _MD_EMPHASIS_RE.sub(r"\2", text)
    return text.strip()


@component
class _MarkdownASTSplitter:
    """Top-level markdown blocks → one atom each. Fenced code blocks are
    never split. ``meta["raw_content"]`` is the rendered markdown source of
    the block (markers preserved).
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        import mistletoe
        from mistletoe.markdown_renderer import MarkdownRenderer

        atoms: list[Document] = []
        for doc in documents:
            content = doc.content or ""
            with MarkdownRenderer() as renderer:
                root = mistletoe.Document(content)
                for tok in root.children:
                    type_name = type(tok).__name__
                    if type_name not in _MD_BLOCK_TYPES:
                        continue  # BlankLine / etc.
                    raw = renderer.render(tok)
                    if not raw.strip():
                        continue
                    text = _md_plain(raw, type_name)
                    if not text:
                        text = raw
                    atoms.append(Document(content=text, meta={**doc.meta, "raw_content": raw}))
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _HtmlASTSplitter (T2v.34/35)
# ---------------------------------------------------------------------------


_HTML_DROP_TAGS = ("script", "style", "nav", "aside", "footer", "header")
_HTML_ATOM_SELECTORS = (
    "h1",
    "h2",
    "h3",
    "h4",
    "h5",
    "h6",
    "p",
    "pre",
    "table",
    "blockquote",
)
_HTML_ATOM_TAGSET = frozenset(_HTML_ATOM_SELECTORS)


@component
class _HtmlASTSplitter:
    """Walks HTML DOM. Drops ``<script>/<style>/<nav>/<aside>/<footer>/<header>``
    (when not nested in ``<article>``/``<main>``); emits one atom per
    block-level element (headings, paragraphs, ``<pre>``, ``<table>``,
    ``<blockquote>``). ``meta["raw_content"]`` is the serialized outer HTML.
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        from selectolax.parser import HTMLParser

        atoms: list[Document] = []
        for doc in documents:
            content = doc.content or ""
            tree = HTMLParser(content)
            self._strip_boilerplate(tree)
            for sel in _HTML_ATOM_SELECTORS:
                for node in tree.css(sel):
                    if self._has_atom_ancestor(node):
                        continue
                    if node.tag == "pre":
                        # Preserve significant whitespace / newlines for code
                        # and pre-formatted blocks. `separator=" "` would
                        # collapse line breaks between text nodes.
                        text = node.text(deep=True)
                    else:
                        text = node.text(deep=True, separator=" ", strip=True)
                    raw = node.html or text
                    if not text.strip():
                        continue
                    atoms.append(
                        Document(
                            content=text,
                            meta={**doc.meta, "raw_content": raw},
                        )
                    )
        return {"documents": atoms}

    @staticmethod
    def _strip_boilerplate(tree: Any) -> None:
        for tag in _HTML_DROP_TAGS:
            for node in tree.css(tag):
                # Keep when nested in an explicit content region.
                anc = node.parent
                inside_content = False
                while anc is not None:
                    if anc.tag in ("article", "main"):
                        inside_content = True
                        break
                    anc = anc.parent
                if not inside_content:
                    node.decompose()

    @staticmethod
    def _has_atom_ancestor(node: Any) -> bool:
        anc = node.parent
        while anc is not None:
            if anc.tag in _HTML_ATOM_TAGSET:
                return True
            anc = anc.parent
        return False


# ---------------------------------------------------------------------------
# _DocxASTSplitter
# ---------------------------------------------------------------------------


def _table_to_markdown(table: Any) -> tuple[str, str]:
    """Render a python-docx Table as (plain_text, markdown_pipe_table).

    Returns both representations in one pass to avoid iterating rows twice.
    """

    def _clean(t: str) -> str:
        return t.replace("|", "\\|").replace("\n", " ")

    rows = [[_clean(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return "", ""
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
    md = "\n".join(filter(None, [header, sep, body]))
    plain = " ".join(cell for row in rows for cell in row if cell.strip())
    return plain, md


@component
class _DocxASTSplitter:
    """DOCX binary → one atom per paragraph / table.

    Reads bytes from ``meta["raw_bytes"]``.  Each heading atom's
    ``content`` is the heading text (no ``#`` markers); each paragraph
    atom's ``content`` is the paragraph text; each table atom's ``content``
    is all cell text joined by spaces and ``meta["raw_content"]`` is the
    Markdown pipe-table representation.
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        from docx import Document as DocxDocument
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        from ragent.security.archive_guard import assert_safe_zip

        atoms: list[Document] = []
        for doc in documents:
            base_meta = {k: v for k, v in doc.meta.items() if k != "raw_bytes"}
            raw_bytes: bytes = doc.meta.get("raw_bytes") or b""
            assert_safe_zip(raw_bytes)
            docx = DocxDocument(io.BytesIO(raw_bytes))
            for block in docx.element.body:
                tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
                if tag == "p":
                    para = Paragraph(block, docx)
                    text = para.text.strip()
                    if not text:
                        continue
                    atoms.append(Document(content=text, meta={**base_meta, "raw_content": text}))
                elif tag == "tbl":
                    table = Table(block, docx)
                    plain, raw_md = _table_to_markdown(table)
                    if not raw_md.strip():
                        continue
                    atoms.append(Document(content=plain, meta={**base_meta, "raw_content": raw_md}))
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _PptxASTSplitter
# ---------------------------------------------------------------------------


@component
class _PptxASTSplitter:
    """PPTX binary → one atom per slide.

    Reads bytes from ``meta["raw_bytes"]``.  Each slide that contains at
    least one non-empty text shape produces one atom.  All text frames on
    the slide are joined with newlines for ``content`` and
    ``meta["raw_content"]``.  ``meta["slide_number"]`` carries the
    1-based slide index.
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        from pptx import Presentation

        from ragent.security.archive_guard import assert_safe_zip

        atoms: list[Document] = []
        for doc in documents:
            base_meta = {k: v for k, v in doc.meta.items() if k != "raw_bytes"}
            raw_bytes: bytes = doc.meta.get("raw_bytes") or b""
            assert_safe_zip(raw_bytes)
            prs = Presentation(io.BytesIO(raw_bytes))
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.type in _PPTX_SKIP_PH:
                        continue
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                texts.append(line)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                cell_text = cell.text_frame.text.strip()
                                if cell_text:
                                    texts.append(cell_text)
                if not texts:
                    continue
                combined = "\n".join(texts)
                atoms.append(
                    Document(
                        content=combined,
                        meta={**base_meta, "raw_content": combined, "slide_number": idx},
                    )
                )
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _PdfASTSplitter
# ---------------------------------------------------------------------------


@component
class _PdfASTSplitter:
    """PDF binary → markdown atoms via pymupdf4llm (RapidOCR auto-selected for image pages).

    Per page: ``pymupdf4llm.to_markdown(pdf, pages=[i], use_ocr=True)`` → markdown string
    → ``_MarkdownASTSplitter`` → structured atoms (headings, paragraphs, tables).
    Empty pages produce no atoms. ``meta["page_number"]`` carries the 1-based index.

    ``fitz.TOOLS.store_shrink(100)`` after each page evicts MuPDF's 256 MB LRU cache,
    bounding peak RSS to one page's worth of intermediate data at a time.
    """

    def __init__(self) -> None:
        self._md_splitter = _MarkdownASTSplitter()

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        import fitz

        from ragent.security.archive_guard import assert_safe_pdf_page_count

        atoms: list[Document] = []
        for doc in documents:
            if not (raw_bytes := doc.meta.get("raw_bytes")):
                continue
            base_meta = {k: v for k, v in doc.meta.items() if k != "raw_bytes"}
            margins = (0, INGEST_PDF_MARGIN_PTS, 0, INGEST_PDF_MARGIN_PTS)
            with fitz.open(stream=raw_bytes, filetype="pdf") as pdf:
                assert_safe_pdf_page_count(pdf.page_count, max_pages=INGEST_MAX_PDF_PAGES)
                for page_idx in range(pdf.page_count):
                    try:
                        md = pymupdf4llm.to_markdown(
                            pdf, pages=[page_idx], use_ocr=True, margins=margins
                        )
                    except Exception:
                        _logger.warning(
                            "pdf_to_markdown_fallback", page=page_idx + 1, exc_info=True
                        )
                        md = pdf[page_idx].get_text("text").strip()
                    fitz.TOOLS.store_shrink(100)
                    if not md.strip():
                        continue
                    page_doc = Document(content=md, meta={**base_meta, "page_number": page_idx + 1})
                    atoms.extend(self._md_splitter.run([page_doc])["documents"])
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _MimeAwareSplitter (T2v.38/39 — replaces FileTypeRouter+joiner+3-splitters)
# ---------------------------------------------------------------------------

_SPLITTER_LABEL: dict[str, str] = {
    "text/plain": "plain",
    "text/markdown": "markdown",
    "text/html": "html",
    IngestMime.DOCX: "docx",
    IngestMime.PPTX: "pptx",
    IngestMime.PDF: "pdf",
}


@component
class _MimeAwareSplitter:
    """Routes Documents to the right splitter based on ``meta["mime_type"]``.

    Single component (not a Haystack ``FileTypeRouter`` + ``DocumentJoiner``
    pair) because Haystack's stock router routes ``ByteStream`` / ``Path``,
    not ``Document``. The plan graph and this implementation are equivalent:
    one fan-in, one fan-out, mime-driven dispatch, unknown → fail.
    """

    def __init__(self) -> None:
        # split_length is in passages; we treat the whole text as one passage
        # and let _BudgetChunker handle sizing.
        self._plain = DocumentSplitter(split_by="passage", split_length=1, split_overlap=0)
        self._plain.warm_up()
        self._md = _MarkdownASTSplitter()
        self._html = _HtmlASTSplitter()
        self._docx = _DocxASTSplitter()
        self._pptx = _PptxASTSplitter()
        self._pdf = _PdfASTSplitter()

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        atoms: list[Document] = []
        for doc in documents:
            mime = doc.meta.get("mime_type") or "text/plain"
            if mime not in _SPLITTER_LABEL:
                raise IngestStepError(
                    f"unroutable mime: {mime!r}", error_code=TaskErrorCode.PIPELINE_UNROUTABLE
                )
            structlog.contextvars.bind_contextvars(splitter=_SPLITTER_LABEL[mime])
            if mime == "text/plain":
                out = self._plain.run([doc])["documents"]
                # DocumentSplitter doesn't set raw_content; default to its content.
                for a in out:
                    a.meta.setdefault("raw_content", a.content or "")
                    a.meta.setdefault("mime_type", mime)
            elif mime == "text/markdown":
                out = self._md.run([doc])["documents"]
            elif mime == "text/html":
                out = self._html.run([doc])["documents"]
            elif mime == IngestMime.DOCX:
                out = self._docx.run([doc])["documents"]
            elif mime == IngestMime.PPTX:
                out = self._pptx.run([doc])["documents"]
            elif mime == IngestMime.PDF:
                out = self._pdf.run([doc])["documents"]
            atoms.extend(out)
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _BudgetChunker (T2v.36/37 — replaces _CharBudgetChunker)
# ---------------------------------------------------------------------------


@component
class _BudgetChunker:
    """Mime-agnostic budget chunker.

    Greedy-packs atoms into chunks ≤ ``CHUNK_TARGET_CHARS`` joined by
    newlines. Atoms longer than ``CHUNK_MAX_CHARS`` are hard-split with
    ``CHUNK_OVERLAP_CHARS`` overlap. Each output ``Document`` carries:
    - ``content``: packed normalized text
    - ``meta["raw_content"]``: concatenation of source atoms' raw slices
    - ``meta["split_id"]``: zero-based per-document index
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        result: list[Document] = []
        # Group by document_id so split_id resets per source document.
        groups: dict[Any, list[Document]] = {}
        order: list[Any] = []
        for d in documents:
            key = d.meta.get("document_id")
            if key not in groups:
                groups[key] = []
                order.append(key)
            groups[key].append(d)

        for doc_id in order:
            atoms = groups[doc_id]
            if not atoms:
                continue
            base_meta = {**atoms[0].meta}
            base_meta.pop("raw_content", None)
            chunks = _pack_atoms(atoms)
            for i, (text, raw) in enumerate(chunks):
                meta = {**base_meta, "split_id": i, "raw_content": raw}
                result.append(Document(content=text, meta=meta))
        return {"documents": result}


def _pack_atoms(atoms: list[Document]) -> list[tuple[str, str]]:
    """Pack atoms into chunks ≤ ``CHUNK_TARGET_CHARS``.

    Text overlap (``CHUNK_OVERLAP_CHARS`` chars carried tail→head between
    adjacent chunks) is for retrieval continuity. ``raw_content`` overlap
    is intentionally NOT carried at character granularity — naive
    ``raw[-overlap:]`` slicing cuts through HTML tags / Markdown
    markers and produces malformed fragments. Instead, ``raw_content`` is
    always atom-aligned: each chunk's raw is the concatenation of the
    raws of the atoms whose normalized text appears in that chunk
    (excluding the carried-overlap prefix).
    """
    target = CHUNK_TARGET_CHARS
    max_chars = CHUNK_MAX_CHARS
    overlap = CHUNK_OVERLAP_CHARS
    chunks: list[tuple[str, str]] = []
    buf_text = ""
    buf_raw = ""

    def flush(carry: bool = True) -> None:
        nonlocal buf_text, buf_raw
        if not buf_text:
            return
        chunks.append((buf_text, buf_raw))
        if carry and overlap > 0 and len(buf_text) > overlap:
            # Carry text overlap for retrieval continuity. Reset raw to ""
            # so the next chunk's raw starts fresh at the next atom
            # boundary — keeps raw_content well-formed for downstream
            # citation rendering.
            buf_text = buf_text[-overlap:]
            buf_raw = ""
        else:
            buf_text = ""
            buf_raw = ""

    for atom in atoms:
        text = atom.content or ""
        raw = atom.meta.get("raw_content") or text
        if not text:
            continue
        if len(text) > max_chars:
            flush(carry=False)
            step = max(1, target - overlap)
            start = 0
            pieces = 0
            while start < len(text):
                end = min(start + target, len(text))
                piece = text[start:end]
                # When raw and text are length-aligned (e.g. plain-text
                # atoms where raw == text), slice raw to match. Otherwise
                # (HTML/markdown atoms whose raw is much larger than the
                # normalized text) fall back to the text piece itself —
                # duplicating the full raw across every hard-split piece
                # would multiply ES storage and LLM tokens by N.
                raw_piece = raw[start:end] if len(raw) == len(text) else piece
                chunks.append((piece, raw_piece))
                pieces += 1
                if pieces > CHUNK_MAX_PIECES_PER_ATOM:
                    raise IngestStepError(
                        f"atom exceeded {CHUNK_MAX_PIECES_PER_ATOM} hard-split pieces "
                        f"(target={target}, overlap={overlap}, atom_len={len(text)})",
                        error_code=TaskErrorCode.CHUNK_BUDGET_EXCEEDED,
                    )
                if end == len(text):
                    break
                start += step
            continue
        sep = "\n" if buf_text else ""
        if buf_text and len(buf_text) + len(sep) + len(text) > target:
            flush(carry=True)
            sep = "\n" if buf_text else ""
        buf_text = buf_text + sep + text
        # Raw is joined only when buf_raw is non-empty (i.e. previous
        # atom contributed); after a flush(carry=True) buf_raw is reset
        # to "" so this chunk's raw begins at the current atom boundary.
        raw_sep = "\n" if buf_raw else ""
        buf_raw = buf_raw + raw_sep + raw
    if buf_text:
        chunks.append((buf_text, buf_raw))
    return chunks


# ---------------------------------------------------------------------------
# DocumentEmbedder — multi-model dual-write (B50 T-EM.15)
# ---------------------------------------------------------------------------


@component
class DocumentEmbedder:
    """Wraps the project's external EmbeddingClient as a Haystack component.

    Two construction modes:

    - **Legacy single-model**: ``DocumentEmbedder(client)`` — embeds every
      chunk with the given client and stores the vector on ``doc.embedding``.
      Kept for tests and the pre-T-EM ingest path.

    - **Registry mode (B50 T-EM-R.6)**: ``DocumentEmbedder(registry, embed_callable, es_client)``
      — reads ``registry.write_models()`` on every ``run()``; bulk-writes to
      ``registry.stable_index`` (and ``registry.candidate_index`` during
      CANDIDATE/CUTOVER) using field ``"embedding"`` in each index. Returns
      ``{"documents": []}`` — the embedder is the sole ES writer, no downstream
      DocumentWriter needed. Empty ``write_models()`` raises ``RuntimeError``.
    """

    def __init__(
        self,
        client: Any = None,
        *,
        registry: Any = None,
        embed_callable: Any = None,
        es_client: Any = None,
    ) -> None:
        if registry is not None:
            if embed_callable is None:
                raise ValueError("registry mode requires embed_callable")
            if es_client is None:
                raise ValueError("registry mode requires es_client")
            self._mode = "dual"
            self._registry = registry
            self._embed = embed_callable
            self._es = es_client
            self._client = None
        else:
            self._mode = "legacy"
            self._client = client
            self._registry = None
            self._embed = None
            self._es = None

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        if not documents:
            return {"documents": []}
        if self._mode == "legacy":
            return self._run_legacy(documents)
        return self._run_dual(documents)

    def _run_legacy(self, documents: list[Document]) -> dict:
        texts = [d.content or "" for d in documents]
        embeddings = self._client.embed(texts)
        out = [
            dataclasses.replace(d, embedding=e) for d, e in zip(documents, embeddings, strict=True)
        ]
        return {"documents": out}

    def _run_dual(self, documents: list[Document]) -> dict:
        models = list(self._registry.write_models())
        if not models:
            raise RuntimeError(
                "ActiveModelRegistry returned no write_models — refusing to emit unindexed chunks"
            )
        texts = [d.content or "" for d in documents]

        if len(models) == 1:
            results = [self._embed(models[0], texts)]
        else:
            with ThreadPoolExecutor(max_workers=len(models)) as pool:
                results = list(pool.map(lambda m: self._embed(m, texts), models))

        index_names = [self._registry.stable_index]
        if len(models) > 1:
            candidate_idx = self._registry.candidate_index
            if candidate_idx is None:
                raise RuntimeError("write_models() returned 2 models but candidate_index is None")
            index_names.append(candidate_idx)

        for vectors, index_name in zip(results, index_names, strict=True):
            ops: list[dict] = []
            op_docs: list[tuple[Document, list[float]]] = []
            for doc, vec in zip(documents, vectors, strict=True):
                ops.append({"index": {"_id": doc.id}})
                # meta is written first so that "embedding" and "content" always win.
                body: dict = {**(doc.meta or {}), "embedding": vec}
                if doc.content is not None:
                    body["content"] = doc.content
                ops.append(body)
                op_docs.append((doc, vec))
            response = self._es.bulk(index=index_name, operations=ops)
            self._handle_bulk_response(response, op_docs, index_name)

        return {"documents": []}

    def _handle_bulk_response(
        self,
        response: dict,
        op_docs: list[tuple[Document, list[float]]],
        index_name: str,
    ) -> None:
        """Check bulk response for partial failures and retry failed items."""
        if not response.get("errors"):
            return

        items = response.get("items", [])
        retry_op_docs: list[tuple[Document, list[float]]] = []
        for item, (doc, vec) in zip(items, op_docs, strict=True):
            action = item.get("index", {})
            status = action.get("status", 200)
            if status >= 400:
                _logger.warning(
                    "es.bulk_partial_failure",
                    index=index_name,
                    doc_id=action.get("_id"),
                    status=status,
                    error=action.get("error"),
                )
                retry_op_docs.append((doc, vec))

        if not retry_op_docs:
            return

        retry_ops: list[dict] = []
        for doc, vec in retry_op_docs:
            retry_ops.append({"index": {"_id": doc.id}})
            body: dict = {**(doc.meta or {}), "embedding": vec}
            if doc.content is not None:
                body["content"] = doc.content
            retry_ops.append(body)
        retry_response = self._es.bulk(index=index_name, operations=retry_ops)
        if retry_response.get("errors"):
            _logger.error(
                "es.bulk_retry_partial_failure",
                index=index_name,
                failed_ids=[
                    item.get("index", {}).get("_id")
                    for item in retry_response.get("items", [])
                    if item.get("index", {}).get("status", 200) >= 400
                ],
            )


# ---------------------------------------------------------------------------
# build_ingest_pipeline — v2 graph
# ---------------------------------------------------------------------------


def build_ingest_pipeline(embedder: Any) -> Pipeline:
    """V2 ingest pipeline.

    Run input shape::

        {
            "loader": {
                "content": str,
                "mime_type": "text/plain"|"text/markdown"|"text/html",
                "document_id": str,
                "source_url": str | None,
                "source_title": str | None,
                ...
            }
        }

    Retry idempotency relies on ``_op_type: index`` (overwrite-by-id) in the
    embedder's bulk writes and deterministic chunk IDs (Haystack hashes
    content+meta). Chunks live only in ES; the v1 ``chunks`` DB table and
    ``ChunkRepository`` were dropped in C6.
    """
    pipeline = Pipeline()

    def _add(name: str, component: Any, *, step: str, error_code: str | None = None) -> None:
        kwargs: dict = {"namespace": "ingest", "step": step}
        if error_code is not None:
            kwargs["error_code"] = error_code
        pipeline.add_component(name, wrap_pipeline_component(component, **kwargs))

    _add("loader", _TextLoader(), step="load")
    _add(
        "splitter",
        _MimeAwareSplitter(),
        step="split",
        error_code=TaskErrorCode.PIPELINE_UNROUTABLE,
    )
    _add("chunker", _BudgetChunker(), step="chunker")
    _add("embedder", embedder, step="embedder", error_code=TaskErrorCode.EMBEDDER_ERROR)

    pipeline.connect("loader.documents", "splitter.documents")
    pipeline.connect("splitter.documents", "chunker.documents")
    pipeline.connect("chunker.documents", "embedder.documents")

    return pipeline
