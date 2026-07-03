"""Splitter components for the v2 ingest pipeline."""

from __future__ import annotations

import csv
import io
import re
from typing import Any

import pymupdf4llm
import structlog
from haystack.components.preprocessors import DocumentSplitter
from haystack.core.component import component
from haystack.dataclasses import Document

from ragent.errors.codes import TaskErrorCode
from ragent.pipelines.observability import IngestStepError
from ragent.schemas.ingest import IngestMime
from ragent.security.archive_guard import INGEST_MAX_PDF_PAGES
from ragent.utility.env import bool_env, float_env, int_env

_logger = structlog.get_logger(__name__)

# PDF page margin in points (1 pt ≈ 0.35 mm). Header/footer zones at the
# top and bottom of each page are excluded from extraction when > 0.
INGEST_PDF_MARGIN_PTS = float_env("INGEST_PDF_MARGIN_PTS", 0.0)

# OCR controls — all default to off / conservative values so text-only PDFs
# complete in seconds without any OCR inference.
_PDF_USE_OCR: bool = bool_env("INGEST_PDF_USE_OCR", False)
_PDF_OCR_CHAR_THRESHOLD: int = int_env("INGEST_PDF_OCR_CHAR_THRESHOLD", 50)
_PDF_OCR_MAX_SCANNED_PAGES: int = int_env("INGEST_PDF_OCR_MAX_SCANNED_PAGES", 10)
# 150 DPI gives the same detection quality as 300 (RapidOCR Det resizes to
# 736px anyway) while making get_pixmap() ~4x faster.
_PDF_OCR_DPI: int = int_env("INGEST_PDF_OCR_DPI", 150)
_PDF_PROGRESS_LOG_INTERVAL: int = 5

# PPTX placeholder types to exclude (header=14, footer=15, date=16, slide_number=13).
# Integer values used so python-pptx stays a lazy import inside run().
_PPTX_SKIP_PH: frozenset[int] = frozenset({13, 14, 15, 16})

# ---------------------------------------------------------------------------
# _MarkdownASTSplitter (T2v.32/33)
# ---------------------------------------------------------------------------


_MD_BLOCK_TYPES = (
    "Heading",
    "Paragraph",
    "CodeFence",
    "List",
    "Table",
    "Quote",
    "ThematicBreak",
    "HtmlBlock",
)

# Strip common markdown markers from a rendered block so the `content`
# field embeds + BM25-indexes prose text rather than syntax noise. Fenced
# code (``` markers) and inline code spans (`x`) are unwrapped to their
# inner text. Heading hashes, bullet/quote prefixes, and emphasis markers
# are dropped. The original markup is always preserved in `raw_content`.
_MD_FENCE_RE = re.compile(r"^```[^\n]*\n(.*?)(?:\n```\s*)?$", re.DOTALL)
_MD_INLINE_CODE_RE = re.compile(r"`+([^`]+)`+")
_MD_HEADING_RE = re.compile(r"^#{1,6}\s+", re.MULTILINE)
_MD_LIST_PREFIX_RE = re.compile(r"^\s*(?:[-*+]|\d+\.)\s+", re.MULTILINE)
_MD_QUOTE_PREFIX_RE = re.compile(r"^>\s?", re.MULTILINE)
_MD_EMPHASIS_RE = re.compile(r"(\*{1,3}|_{1,3})(\S(?:.*?\S)?)\1")


def _md_plain(raw: str, type_name: str) -> str:
    """Reduce a rendered markdown block to its prose content."""
    if type_name == "CodeFence":
        m = _MD_FENCE_RE.match(raw.rstrip())
        return (m.group(1) if m else raw).rstrip()
    text = raw
    if type_name == "Heading":
        text = _MD_HEADING_RE.sub("", text)
    if type_name in ("List", "Quote"):
        text = _MD_LIST_PREFIX_RE.sub("", text)
        text = _MD_QUOTE_PREFIX_RE.sub("", text)
    text = _MD_INLINE_CODE_RE.sub(r"\1", text)
    text = _MD_EMPHASIS_RE.sub(r"\2", text)
    return text.strip()


@component
class _MarkdownASTSplitter:
    """Top-level markdown blocks → one atom each. Fenced code blocks are
    never split. ``meta["raw_content"]`` is the rendered markdown source of
    the block (markers preserved).
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        import mistletoe
        from mistletoe.markdown_renderer import MarkdownRenderer

        atoms: list[Document] = []
        for doc in documents:
            content = doc.content or ""
            with MarkdownRenderer() as renderer:
                root = mistletoe.Document(content)
                for tok in root.children:
                    type_name = type(tok).__name__
                    if type_name not in _MD_BLOCK_TYPES:
                        continue  # BlankLine / etc.
                    raw = renderer.render(tok)
                    if not raw.strip():
                        continue
                    text = _md_plain(raw, type_name)
                    if not text:
                        text = raw
                    atoms.append(Document(content=text, meta={**doc.meta, "raw_content": raw}))
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _HtmlASTSplitter (T2v.34/35)
# ---------------------------------------------------------------------------


_HTML_DROP_TAGS = ("script", "style", "nav", "aside", "footer", "header")
_HTML_ATOM_SELECTORS = (
    "h1",
    "h2",
    "h3",
    "h4",
    "h5",
    "h6",
    "p",
    "pre",
    "table",
    "blockquote",
)
_HTML_ATOM_TAGSET = frozenset(_HTML_ATOM_SELECTORS)


@component
class _HtmlASTSplitter:
    """Walks HTML DOM. Drops ``<script>/<style>/<nav>/<aside>/<footer>/<header>``
    (when not nested in ``<article>``/``<main>``); emits one atom per
    block-level element (headings, paragraphs, ``<pre>``, ``<table>``,
    ``<blockquote>``). ``meta["raw_content"]`` is the serialized outer HTML.
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        from selectolax.parser import HTMLParser

        atoms: list[Document] = []
        for doc in documents:
            content = doc.content or ""
            tree = HTMLParser(content)
            self._strip_boilerplate(tree)
            for sel in _HTML_ATOM_SELECTORS:
                for node in tree.css(sel):
                    if self._has_atom_ancestor(node):
                        continue
                    if node.tag == "pre":
                        # Preserve significant whitespace / newlines for code
                        # and pre-formatted blocks. `separator=" "` would
                        # collapse line breaks between text nodes.
                        text = node.text(deep=True)
                    else:
                        text = node.text(deep=True, separator=" ", strip=True)
                    raw = node.html or text
                    if not text.strip():
                        continue
                    atoms.append(
                        Document(
                            content=text,
                            meta={**doc.meta, "raw_content": raw},
                        )
                    )
        return {"documents": atoms}

    @staticmethod
    def _strip_boilerplate(tree: Any) -> None:
        for tag in _HTML_DROP_TAGS:
            for node in tree.css(tag):
                # Keep when nested in an explicit content region.
                anc = node.parent
                inside_content = False
                while anc is not None:
                    if anc.tag in ("article", "main"):
                        inside_content = True
                        break
                    anc = anc.parent
                if not inside_content:
                    node.decompose()

    @staticmethod
    def _has_atom_ancestor(node: Any) -> bool:
        anc = node.parent
        while anc is not None:
            if anc.tag in _HTML_ATOM_TAGSET:
                return True
            anc = anc.parent
        return False


# ---------------------------------------------------------------------------
# _DocxASTSplitter
# ---------------------------------------------------------------------------


_DOCX_HEADING_RE = re.compile(r"^Heading\s*(\d+)$", re.IGNORECASE)


def _docx_heading_level(style_name: str) -> int | None:
    """Map a python-docx paragraph style name to a markdown heading level
    (1-6), or None if the style is not a heading. "Title" is treated as
    level 1; "Heading N" is clamped to 6 (markdown's maximum)."""
    name = style_name.strip()
    if name.lower() == "title":
        return 1
    if m := _DOCX_HEADING_RE.match(name):
        return min(int(m.group(1)), 6)
    return None


def _table_to_markdown(table: Any) -> tuple[str, str]:
    """Render a python-docx Table as (plain_text, markdown_pipe_table).

    Returns both representations in one pass to avoid iterating rows twice.
    """

    def _clean(t: str) -> str:
        return t.replace("|", "\\|").replace("\n", " ")

    rows = [[_clean(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return "", ""
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
    md = "\n".join(filter(None, [header, sep, body]))
    plain = " ".join(cell for row in rows for cell in row if cell.strip())
    return plain, md


@component
class _DocxASTSplitter:
    """DOCX binary → one atom per paragraph / table.

    Reads bytes from ``meta["raw_bytes"]``.  Each heading atom's
    ``content`` is the heading text (no ``#`` markers) while
    ``meta["raw_content"]`` carries the markdown ``#``-prefixed form (level
    derived from the paragraph's "Heading N"/"Title" style); each plain
    paragraph atom's ``content``/``raw_content`` is just the paragraph text;
    each table atom's ``content`` is all cell text joined by spaces and
    ``meta["raw_content"]`` is the Markdown pipe-table representation.
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        from docx import Document as DocxDocument
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        from ragent.security.archive_guard import assert_safe_zip

        atoms: list[Document] = []
        for doc in documents:
            base_meta = {k: v for k, v in doc.meta.items() if k != "raw_bytes"}
            raw_bytes: bytes = doc.meta.get("raw_bytes") or b""
            assert_safe_zip(raw_bytes)
            docx = DocxDocument(io.BytesIO(raw_bytes))
            for block in docx.element.body:
                tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
                if tag == "p":
                    para = Paragraph(block, docx)
                    text = para.text.strip()
                    if not text:
                        continue
                    level = _docx_heading_level(para.style.name if para.style else "")
                    raw_content = f"{'#' * level} {text}" if level else text
                    atoms.append(
                        Document(content=text, meta={**base_meta, "raw_content": raw_content})
                    )
                elif tag == "tbl":
                    table = Table(block, docx)
                    plain, raw_md = _table_to_markdown(table)
                    if not raw_md.strip():
                        continue
                    atoms.append(Document(content=plain, meta={**base_meta, "raw_content": raw_md}))
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _PptxASTSplitter
# ---------------------------------------------------------------------------


def _pptx_atom(content: str, raw_content: str, base_meta: dict, slide_number: int) -> Document:
    return Document(
        content=content,
        meta={**base_meta, "raw_content": raw_content, "slide_number": slide_number},
    )


@component
class _PptxASTSplitter:
    """PPTX binary → one heading atom + one body atom per slide.

    Reads bytes from ``meta["raw_bytes"]``.  A slide's title placeholder
    (``slide.shapes.title``), if present and non-empty, becomes its own
    heading atom — ``content`` is the title text, ``meta["raw_content"]``
    is the markdown ``# <title>`` form. The remaining (non-title,
    non-skipped) text frames on the slide are joined with newlines into a
    second, body atom. ``meta["slide_number"]`` carries the 1-based slide
    index on both atoms. A slide with neither a title nor any other text
    produces no atom.
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        from pptx import Presentation

        from ragent.security.archive_guard import assert_safe_zip

        atoms: list[Document] = []
        for doc in documents:
            base_meta = {k: v for k, v in doc.meta.items() if k != "raw_bytes"}
            raw_bytes: bytes = doc.meta.get("raw_bytes") or b""
            assert_safe_zip(raw_bytes)
            prs = Presentation(io.BytesIO(raw_bytes))
            for idx, slide in enumerate(prs.slides, start=1):
                title_shape = slide.shapes.title
                title_text = ""
                if title_shape is not None and title_shape.has_text_frame:
                    title_text = title_shape.text_frame.text.strip()

                texts = []
                for shape in slide.shapes:
                    if title_shape is not None and shape.shape_id == title_shape.shape_id:
                        continue
                    if shape.is_placeholder and shape.placeholder_format.type in _PPTX_SKIP_PH:
                        continue
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                texts.append(line)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                cell_text = cell.text_frame.text.strip()
                                if cell_text:
                                    texts.append(cell_text)

                if title_text:
                    atoms.append(_pptx_atom(title_text, f"# {title_text}", base_meta, idx))
                if texts:
                    combined = "\n".join(texts)
                    atoms.append(_pptx_atom(combined, combined, base_meta, idx))
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _PdfASTSplitter
# ---------------------------------------------------------------------------


@component
class _PdfASTSplitter:
    """PDF binary → markdown atoms via pymupdf4llm.

    Per page: ``pymupdf4llm.to_markdown(pdf, pages=[i], use_ocr=<per-page>)`` →
    markdown string → ``_MarkdownASTSplitter`` → structured atoms.
    Empty pages produce no atoms. ``meta["page_number"]`` carries the 1-based index.

    OCR behaviour (controlled by env vars):
    - ``INGEST_PDF_USE_OCR=false`` (default): no OCR; text-layer extraction only.
    - ``INGEST_PDF_USE_OCR=true``: pre-scans all pages cheaply via ``get_text()``;
      pages with fewer than ``INGEST_PDF_OCR_CHAR_THRESHOLD`` chars are flagged as
      scanned and OCR'd individually.  If the scanned-page count exceeds
      ``INGEST_PDF_OCR_MAX_SCANNED_PAGES`` the task is rejected immediately
      (``PdfTooManyScannedPagesError``) so the worker can write a typed FAILED row.

    ``fitz.TOOLS.store_shrink(100)`` after each page evicts MuPDF's 256 MB LRU
    cache, bounding peak RSS to one page's worth of intermediate data at a time.
    """

    def __init__(self) -> None:
        self._md_splitter = _MarkdownASTSplitter()

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        import fitz

        from ragent.security.archive_guard import (
            PdfTooManyScannedPagesError,
            assert_safe_pdf_page_count,
        )

        atoms: list[Document] = []
        for doc in documents:
            if not (raw_bytes := doc.meta.get("raw_bytes")):
                continue
            base_meta = {k: v for k, v in doc.meta.items() if k != "raw_bytes"}
            margins = (0, INGEST_PDF_MARGIN_PTS, 0, INGEST_PDF_MARGIN_PTS)
            with fitz.open(stream=raw_bytes, filetype="pdf") as pdf:
                assert_safe_pdf_page_count(pdf.page_count, max_pages=INGEST_MAX_PDF_PAGES)

                # Cheap pre-scan: get_text("text") avoids pixmap rendering
                # (~0.01 s/page). Stops as soon as the cap is exceeded so large
                # fully-scanned PDFs are rejected after cap+1 pages, not all N.
                if _PDF_USE_OCR:
                    scanned_pages: set[int] = set()
                    for i in range(pdf.page_count):
                        if len(pdf[i].get_text("text").strip()) < _PDF_OCR_CHAR_THRESHOLD:
                            scanned_pages.add(i)
                            if len(scanned_pages) > _PDF_OCR_MAX_SCANNED_PAGES:
                                raise PdfTooManyScannedPagesError(
                                    len(scanned_pages), _PDF_OCR_MAX_SCANNED_PAGES
                                )
                    _logger.info(
                        "pdf_ocr_plan",
                        total_pages=pdf.page_count,
                        scanned_pages=len(scanned_pages),
                        ocr_dpi=_PDF_OCR_DPI,
                    )
                else:
                    scanned_pages = set()

                for page_idx in range(pdf.page_count):
                    page_use_ocr = page_idx in scanned_pages
                    if page_idx % _PDF_PROGRESS_LOG_INTERVAL == 0:
                        _logger.debug(
                            "pdf_page_progress",
                            page=page_idx + 1,
                            total=pdf.page_count,
                            ocr=page_use_ocr,
                        )
                    try:
                        md = pymupdf4llm.to_markdown(
                            pdf,
                            pages=[page_idx],
                            use_ocr=page_use_ocr,
                            ocr_dpi=_PDF_OCR_DPI,
                            margins=margins,
                        )
                    except Exception:
                        _logger.warning(
                            "pdf_to_markdown_fallback", page=page_idx + 1, exc_info=True
                        )
                        md = pdf[page_idx].get_text("text").strip()
                    fitz.TOOLS.store_shrink(100)
                    if not md.strip():
                        continue
                    page_doc = Document(content=md, meta={**base_meta, "page_number": page_idx + 1})
                    atoms.extend(self._md_splitter.run([page_doc])["documents"])
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _CsvASTSplitter (T-CAT.9)
# ---------------------------------------------------------------------------


@component
class _CsvASTSplitter:
    """CSV rows → one atom per data row. Header row names become field labels.
    Each atom's content is a `: `-delimited line (e.g., "name: alice, age: 30").
    """

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        atoms: list[Document] = []
        for doc in documents:
            content = doc.content or ""
            reader = csv.DictReader(io.StringIO(content))
            if reader.fieldnames is None:
                continue
            for row in reader:
                parts = [f"{k}: {v}" for k, v in row.items() if k is not None and v]
                if not parts:
                    continue
                row_str = ", ".join(parts)
                atoms.append(
                    Document(
                        content=row_str,
                        meta={**doc.meta},
                    )
                )
        return {"documents": atoms}


# ---------------------------------------------------------------------------
# _MimeAwareSplitter (T2v.38/39 — replaces FileTypeRouter+joiner+3-splitters)
# ---------------------------------------------------------------------------

_SPLITTER_LABEL: dict[str, str] = {
    "text/plain": "plain",
    "text/markdown": "markdown",
    "text/html": "html",
    "text/csv": "csv",
    IngestMime.DOCX: "docx",
    IngestMime.PPTX: "pptx",
    IngestMime.PDF: "pdf",
}


@component
class _MimeAwareSplitter:
    """Routes Documents to the right splitter based on ``meta["mime_type"]``.

    Single component (not a Haystack ``FileTypeRouter`` + ``DocumentJoiner``
    pair) because Haystack's stock router routes ``ByteStream`` / ``Path``,
    not ``Document``. The plan graph and this implementation are equivalent:
    one fan-in, one fan-out, mime-driven dispatch, unknown → fail.
    """

    def __init__(self) -> None:
        # split_length is in passages; we treat the whole text as one passage
        # and let _BudgetChunker handle sizing.
        self._plain = DocumentSplitter(split_by="passage", split_length=1, split_overlap=0)
        self._plain.warm_up()
        self._md = _MarkdownASTSplitter()
        self._html = _HtmlASTSplitter()
        self._csv = _CsvASTSplitter()
        self._docx = _DocxASTSplitter()
        self._pptx = _PptxASTSplitter()
        self._pdf = _PdfASTSplitter()

    @component.output_types(documents=list[Document])
    def run(self, documents: list[Document]) -> dict:
        atoms: list[Document] = []
        for doc in documents:
            mime = doc.meta.get("mime_type") or "text/plain"
            if mime not in _SPLITTER_LABEL:
                raise IngestStepError(
                    f"unroutable mime: {mime!r}", error_code=TaskErrorCode.PIPELINE_UNROUTABLE
                )
            structlog.contextvars.bind_contextvars(splitter=_SPLITTER_LABEL[mime])
            if mime == "text/plain":
                out = self._plain.run([doc])["documents"]
                # DocumentSplitter doesn't set raw_content; default to its content.
                for a in out:
                    a.meta.setdefault("raw_content", a.content or "")
                    a.meta.setdefault("mime_type", mime)
            elif mime == "text/markdown":
                out = self._md.run([doc])["documents"]
            elif mime == "text/html":
                out = self._html.run([doc])["documents"]
            elif mime == "text/csv":
                out = self._csv.run([doc])["documents"]
            elif mime == IngestMime.DOCX:
                out = self._docx.run([doc])["documents"]
            elif mime == IngestMime.PPTX:
                out = self._pptx.run([doc])["documents"]
            elif mime == IngestMime.PDF:
                out = self._pdf.run([doc])["documents"]
            atoms.extend(out)
        return {"documents": atoms}
